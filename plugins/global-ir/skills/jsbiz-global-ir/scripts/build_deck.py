#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
jsbiz-global-ir deck builder.

deck_spec.json -> 16:9 PPTX (10 slide-type library, 5 themes, speaker notes).
Pure python-pptx + Pillow. No Node/npm required.

Usage:
    python build_deck.py spec.json [--out out.pptx]

Spec: see assets/spec.schema.json and assets/spec.example.json.
Coordinates are inches on a 13.333 x 7.5 canvas. Do not change layout
constants casually - they are regression-tested against the reference deck.
"""
import argparse
import copy
import json
import os
import sys
import tempfile

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

FONT = "Arial"
SLIDE_W, SLIDE_H = 13.333, 7.5

# ---------------------------------------------------------------- themes ----
FIXED = {"WHITE": "FFFFFF", "RED": "C0392B", "AMBER": "D68910"}

THEMES = {
    "ocean": dict(DEEP="132B40", HEAD="0E3A5D", SUB="0A7A8C", ACCENT="02C39A",
                  ACCENT_TINT="E4F7F1", LIGHT="EEF6FA", ICE="CADCFC", TXT="1A2B3C",
                  MUT="5A6B7B", BORDER="DDE7EE", DEEP_CARD="1B3A52", DEEP_LINE="2A5372",
                  DEEP_ICON="0F2A3D", STRIPE="F6FAFC"),
    "midnight": dict(DEEP="141A3C", HEAD="1E2761", SUB="3A4FB8", ACCENT="8FB0FF",
                  ACCENT_TINT="EAF0FF", LIGHT="EEF1FB", ICE="CADCFC", TXT="20263F",
                  MUT="5C647F", BORDER="DBE1F0", DEEP_CARD="202A56", DEEP_LINE="32407A",
                  DEEP_ICON="10163A", STRIPE="F5F7FC"),
    "forest": dict(DEEP="1C3320", HEAD="2C5F2D", SUB="4E7B50", ACCENT="97BC62",
                  ACCENT_TINT="EFF5E4", LIGHT="F0F5EE", ICE="D7E8D0", TXT="22301F",
                  MUT="5F6E58", BORDER="DCE7D8", DEEP_CARD="27452B", DEEP_LINE="3B5F40",
                  DEEP_ICON="152718", STRIPE="F6FAF4"),
    "terracotta": dict(DEEP="4E241D", HEAD="8A3A2E", SUB="6E8B7B", ACCENT="D98E4A",
                  ACCENT_TINT="F8ECDD", LIGHT="F5F1E8", ICE="F0D9CE", TXT="33231F",
                  MUT="74655F", BORDER="E7DCD2", DEEP_CARD="63302A", DEEP_LINE="7E453C",
                  DEEP_ICON="3B1B15", STRIPE="FAF6F0"),
    "charcoal": dict(DEEP="23272B", HEAD="2F3B44", SUB="4E6373", ACCENT="E3A72F",
                  ACCENT_TINT="F6EEDC", LIGHT="F2F4F5", ICE="D9DEE3", TXT="22272B",
                  MUT="5F6B75", BORDER="DEE3E7", DEEP_CARD="2E353B", DEEP_LINE="46505A",
                  DEEP_ICON="1A1E22", STRIPE="F7F8F9"),
}

ICON_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "assets", "icons")


# ---------------------------------------------------------------- helpers ---
class Ctx:
    """Build context: theme palette, asset resolution, icon tint cache."""

    def __init__(self, spec, spec_dir):
        theme_name = spec.get("theme", "ocean")
        if theme_name not in THEMES:
            raise SystemExit(f"unknown theme '{theme_name}' (choose from {list(THEMES)})")
        self.pal = dict(THEMES[theme_name])
        self.pal.update(FIXED)
        self.spec_dir = spec_dir
        self.assets_dir = spec.get("assets_dir", "")
        self.tmp = tempfile.mkdtemp(prefix="jsbiz_ir_")
        self._tint_cache = {}

    def color(self, name_or_hex):
        """Resolve a palette role name ('accent') or literal hex ('FF0000')."""
        if not name_or_hex:
            return None
        key = str(name_or_hex)
        if key.upper() in self.pal:
            return self.pal[key.upper()]
        if key in self.pal:
            return self.pal[key]
        return key.lstrip("#").upper()

    def rgb(self, name_or_hex):
        return RGBColor.from_string(self.color(name_or_hex))

    def asset(self, path):
        """Resolve an image path: absolute, or relative to assets_dir / spec dir."""
        if not path:
            return None
        if os.path.isabs(path) and os.path.exists(path):
            return path
        for base in (self.assets_dir, self.spec_dir):
            if base:
                p = os.path.join(base, path)
                if os.path.exists(p):
                    return p
        raise SystemExit(f"asset not found: {path}")

    def icon(self, name, color_role):
        """Tint the bundled white icon PNG to a palette color. Cached per run."""
        hexcol = self.color(color_role)
        key = (name, hexcol)
        if key in self._tint_cache:
            return self._tint_cache[key]
        src = os.path.join(ICON_DIR, f"{name}.png")
        if not os.path.exists(src):
            raise SystemExit(f"unknown icon '{name}' (see assets/icons/)")
        img = Image.open(src).convert("RGBA")
        r, g, b = (int(hexcol[i:i + 2], 16) for i in (0, 2, 4))
        solid = Image.new("RGBA", img.size, (r, g, b, 0))
        solid.putalpha(img.getchannel("A"))
        out = os.path.join(self.tmp, f"{name}_{hexcol}.png")
        solid.save(out)
        self._tint_cache[key] = out
        return out


def _set_fill_alpha(shape, transparency_pct):
    """PowerPoint-style transparency (0=opaque..100) on a solid fill."""
    alpha = int((100 - transparency_pct) * 1000)
    sppr = shape._element.spPr
    srgb = sppr.find(qn("a:solidFill") + "/" + qn("a:srgbClr"))
    if srgb is None:
        srgb = sppr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    el = srgb.makeelement(qn("a:alpha"), {"val": str(alpha)})
    srgb.append(el)


def rect(ctx, slide, x, y, w, h, fill, line=None, line_w=1.0, radius=None,
         transparency=None):
    """Rectangle / rounded rectangle with themed fill. radius in inches."""
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        sp.adjustments[0] = max(0.0, min(0.5, radius / min(w, h)))
    sp.fill.solid()
    sp.fill.fore_color.rgb = ctx.rgb(fill)
    if transparency:
        _set_fill_alpha(sp, transparency)
    if line:
        sp.line.color.rgb = ctx.rgb(line)
        sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def text(ctx, slide, x, y, w, h, content, size=12, color="TXT", bold=False,
         italic=False, align=None, valign=None, line_spacing=None,
         char_spacing=None):
    """Text box. content = str (\\n -> paragraphs) or list of run dicts:
    {text, bold?, italic?, color?, break?}  (break=True starts a new paragraph
    after the run, mirroring pptxgenjs breakLine)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    if isinstance(content, str):
        paras = content.split("\n")
        runsets = [[{"text": p}] for p in paras]
    else:
        runsets, cur = [], []
        for r in content:
            cur.append(r)
            if r.get("break"):
                runsets.append(cur)
                cur = []
        if cur:
            runsets.append(cur)

    for i, runs in enumerate(runsets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        for rspec in runs:
            r = p.add_run()
            r.text = rspec["text"]
            f = r.font
            f.name = FONT
            f.size = Pt(rspec.get("size", size))
            f.bold = rspec.get("bold", bold)
            f.italic = rspec.get("italic", italic)
            f.color.rgb = ctx.rgb(rspec.get("color", color))
            if char_spacing:
                f._rPr.set("spc", str(int(char_spacing * 100)))
    return tb


def image(ctx, slide, path, x, y, w=None, h=None):
    kw = {}
    if w:
        kw["width"] = Inches(w)
    if h:
        kw["height"] = Inches(h)
    return slide.shapes.add_picture(ctx.asset(path), Inches(x), Inches(y), **kw)


def icon_circle(ctx, slide, icon, color_role, x, y, d=0.52, bg="LIGHT"):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid()
    c.fill.fore_color.rgb = ctx.rgb(bg)
    c.line.fill.background()
    c.shadow.inherit = False
    pad = d * 0.26
    slide.shapes.add_picture(ctx.icon(icon, color_role), Inches(x + pad), Inches(y + pad),
                             width=Inches(d - 2 * pad), height=Inches(d - 2 * pad))


def kicker(ctx, slide, num, label, brand, dark=False):
    text(ctx, slide, 0.6, 0.32, 8, 0.32, f"{num}  ·  {label}", size=12,
         bold=True, color="ACCENT", char_spacing=3)
    text(ctx, slide, 11.2, 0.32, 1.55, 0.32, brand, size=10, bold=True,
         color="ICE" if dark else "MUT", align="right", char_spacing=2)


def title(ctx, slide, txt, dark=False, y=0.66, w=12.1):
    text(ctx, slide, 0.58, y, w, 0.85, txt, size=30, bold=True,
         color="WHITE" if dark else "HEAD")


def page_no(ctx, slide, n, total, dark=False):
    text(ctx, slide, 12.35, 7.08, 0.8, 0.3, f"{n} / {total}", size=9,
         color="ICE" if dark else "MUT", align="right")


def _chart_no_title(chart):
    """Force autoTitleDeleted so PowerPoint never invents a chart title."""
    el = chart._chartSpace.find(qn("c:chart"))
    if el.find(qn("c:autoTitleDeleted")) is None:
        atd = el.makeelement(qn("c:autoTitleDeleted"), {"val": "1"})
        el.insert(0, atd)
    else:
        el.find(qn("c:autoTitleDeleted")).set("val", "1")


def _hide_axis(axis):
    """value-axis hidden (c:delete val=1), keeping schema element order."""
    el = axis._element
    d = el.find(qn("c:delete"))
    if d is None:
        d = el.makeelement(qn("c:delete"), {"val": "1"})
        el.find(qn("c:scaling")).addnext(d)
    else:
        d.set("val", "1")


def bar_chart(ctx, slide, x, y, w, h, categories, series, colors, num_fmt="0.0",
              label_size=10, max_val=None, gap=60):
    """Clustered column chart, overlap 100 (so None-padded series stack in place),
    value axis hidden, category axis quiet - the deck's house chart style."""
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y),
                                Inches(w), Inches(h), data)
    ch = gf.chart
    ch.has_legend = False
    _chart_no_title(ch)
    plot = ch.plots[0]
    plot.gap_width = gap
    plot.overlap = 100
    for s, col in zip(ch.series, colors):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = ctx.rgb(col)
        s.format.line.fill.background()
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    dl.number_format = num_fmt
    dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.size = Pt(label_size)
    dl.font.name = FONT
    dl.font.color.rgb = ctx.rgb("HEAD")
    ca = ch.category_axis
    ca.has_major_gridlines = False
    ca.tick_labels.font.size = Pt(10)
    ca.tick_labels.font.name = FONT
    ca.tick_labels.font.color.rgb = ctx.rgb("MUT")
    ca.format.line.color.rgb = ctx.rgb("BORDER")
    va = ch.value_axis
    va.has_major_gridlines = False
    if max_val:
        va.maximum_scale = max_val
        va.minimum_scale = 0
    _hide_axis(va)
    return ch


# ------------------------------------------------------------ slide types ---
def s_cover_dark(ctx, prs, slide, sp, meta):
    if sp.get("bg_image"):
        image(ctx, slide, sp["bg_image"], 0, 2.92, w=SLIDE_W, h=4.58)
        rect(ctx, slide, 0, 2.92, SLIDE_W, 4.58, "DEEP", transparency=28)
    if sp.get("logo"):
        image(ctx, slide, sp["logo"], 0.75, 0.62, w=sp.get("logo_w", 1.95))
    if sp.get("program_line"):
        text(ctx, slide, 0.78, 1.52, 8.5, 0.34, sp["program_line"], size=12,
             color="ICE", char_spacing=1.5)
    runs = []
    for i, ln in enumerate(sp["headline_lines"]):
        runs.append({"text": ln["text"], "color": "ACCENT" if ln.get("accent") else "WHITE",
                     "break": i < len(sp["headline_lines"]) - 1})
    text(ctx, slide, 0.72, 2.25, 8.6, 2.0, runs, size=48, bold=True, line_spacing=56)
    if sp.get("subtitle"):
        text(ctx, slide, 0.75, 4.42, 7.3, 0.95, sp["subtitle"], size=15, color="ICE",
             line_spacing=22)
    if sp.get("product_image"):
        rect(ctx, slide, 8.62, 3.55, 4.05, 2.75, "WHITE", radius=0.14)
        img = Image.open(ctx.asset(sp["product_image"]))
        iw = 3.65
        ih = iw * img.height / img.width
        image(ctx, slide, sp["product_image"], 8.82, 3.55 + (2.35 - ih) / 2 + 0.15, w=iw, h=ih)
        if sp.get("product_caption"):
            text(ctx, slide, 8.62, 5.9, 4.05, 0.3, sp["product_caption"], size=10.5,
                 color="HEAD", bold=True, align="center")
    if sp.get("footer"):
        text(ctx, slide, 0.75, 6.85, 8.5, 0.32, sp["footer"], size=11, color="ICE")


def s_stat_tiles(ctx, prs, slide, sp, meta):
    for i, t in enumerate(sp["tiles"][:4]):
        x = 0.6 + i * 3.11
        rect(ctx, slide, x, 1.72, 2.92, 1.32, "LIGHT", radius=0.1)
        text(ctx, slide, x + 0.18, 1.86, 2.6, 0.6, t["big"], size=30, bold=True, color="SUB")
        text(ctx, slide, x + 0.18, 2.5, 2.6, 0.42, t["small"], size=11.5, color="MUT")
    text(ctx, slide, 0.6, 3.35, 4, 0.35, sp.get("left_heading", ""), size=15,
         bold=True, color="HEAD")
    for i, r in enumerate(sp["left_rows"][:3]):
        y = 3.82 + i * 0.78
        icon_circle(ctx, slide, r["icon"], "SUB", 0.6, y, 0.5)
        text(ctx, slide, 1.28, y - 0.06, 4.62, 0.72, r["text"], size=12.5, valign="middle")
    mc = sp.get("middle_card")
    if mc:
        rect(ctx, slide, 6.2, 3.42, 3.4, 2.9, "WHITE", line="BORDER", radius=0.1)
        text(ctx, slide, 6.4, 3.56, 3.0, 0.3, mc["label"], size=11, bold=True, color="MUT")
        img = Image.open(ctx.asset(mc["image"]))
        iw = 2.8
        ih = min(2.16, iw * img.height / img.width)
        image(ctx, slide, mc["image"], 6.5, 3.95, w=iw, h=ih)
    ch = sp.get("chart")
    if ch:
        text(ctx, slide, 9.95, 3.42, 3.0, 0.3, ch["label"], size=11, bold=True, color="MUT")
        series = [("Actual", [v if v is not None else None for v in ch["actual"]])]
        colors = ["SUB"]
        if ch.get("target"):
            series.append(("Target", ch["target"]))
            colors.append("ACCENT")
        bar_chart(ctx, slide, 9.85, 3.75, 3.0, 2.55, ch["categories"], series, colors,
                  num_fmt=ch.get("number_format", "0.00"),
                  max_val=ch.get("max_val"))
    if sp.get("bottom_line"):
        text(ctx, slide, 0.6, 6.55, 12.1, 0.35, sp["bottom_line"], size=11, color="MUT")


def s_product_hero(ctx, prs, slide, sp, meta):
    img = Image.open(ctx.asset(sp["hero_image"]))
    iw = 5.7
    ih = iw * img.height / img.width
    image(ctx, slide, sp["hero_image"], 7.05, 1.75, w=iw, h=ih)
    if sp.get("hero_caption"):
        text(ctx, slide, 7.05, 1.75 + ih + 0.04, 5.7, 0.3, sp["hero_caption"], size=10.5,
             color="MUT", align="center", italic=True)
    for i, p in enumerate(sp["pillars"][:3]):
        y = 1.85 + i * 1.12
        icon_circle(ctx, slide, p["icon"], "SUB", 0.6, y, 0.56)
        text(ctx, slide, 1.38, y - 0.04, 5.4, 0.34, p["head"], size=15, bold=True, color="HEAD")
        text(ctx, slide, 1.38, y + 0.3, 5.4, 0.66, p["body"], size=12, color="MUT")
    band = sp.get("band")
    if band:
        rect(ctx, slide, 0.6, 5.55, 12.13, 1.35, "LIGHT", radius=0.1)
        tx = 1.85 if band.get("image") else 0.95
        if band.get("image"):
            bimg = Image.open(ctx.asset(band["image"]))
            bh = 1.01
            bw = bh * bimg.width / bimg.height
            image(ctx, slide, band["image"], 0.95, 5.72, w=bw, h=bh)
        text(ctx, slide, tx, 5.72, 5.5, 0.32, band["head"], size=13.5, bold=True, color="HEAD")
        text(ctx, slide, tx, 6.06, 12.45 - tx, 0.7, band["body"], size=12)


def s_icon_rows_photo(ctx, prs, slide, sp, meta):
    for i, r in enumerate(sp["rows"][:3]):
        y = 1.95 + i * 1.24
        icon_circle(ctx, slide, r["icon"], "SUB", 0.6, y, 0.56)
        text(ctx, slide, 1.38, y - 0.04, 6.6, 0.34, r["head"], size=15.5, bold=True, color="HEAD")
        text(ctx, slide, 1.38, y + 0.31, 6.55, 0.72, r["body"], size=12.5, color="MUT")
    if sp.get("photo_main"):
        image(ctx, slide, sp["photo_main"], 8.75, 1.9, w=3.95, h=3.56)
    if sp.get("photo_thumb"):
        image(ctx, slide, sp["photo_thumb"], 8.45, 4.5, w=1.2, h=1.36)
    if sp.get("photo_caption"):
        text(ctx, slide, 9.8, 5.56, 2.9, 0.72, sp["photo_caption"], size=10,
             italic=True, color="MUT")
    if sp.get("closing_runs"):
        rect(ctx, slide, 0.6, 6.35, 12.13, 0.72, "LIGHT", radius=0.09)
        text(ctx, slide, 0.85, 6.42, 11.7, 0.58, sp["closing_runs"], size=12.5,
             valign="middle")


def s_comparison_table(ctx, prs, slide, sp, meta):
    col_x, col_w = [4.35, 7.25, 10.15], 2.75
    rect(ctx, slide, col_x[2] - 0.18, 1.62, col_w + 0.36, 5.45, "ACCENT_TINT",
         line="ACCENT", line_w=1.75, radius=0.12)
    for i, h in enumerate(sp["columns"][:3]):
        if h.get("image"):
            img = Image.open(ctx.asset(h["image"]))
            ih = h.get("img_h", 1.18)
            iw = ih * img.width / img.height
            if iw > col_w - 0.3:
                iw = col_w - 0.3
                ih = iw * img.height / img.width
            image(ctx, slide, h["image"], col_x[i] + (col_w - iw) / 2,
                  1.78 + (1.18 - ih) / 2, w=iw, h=ih)
        text(ctx, slide, col_x[i], 3.0, col_w, 0.3, h["name"], size=13, bold=True,
             color="SUB" if i == 2 else "HEAD", align="center")
        text(ctx, slide, col_x[i], 3.28, col_w, 0.26, h.get("sub", ""), size=10,
             color="MUT", align="center")
    for r, row in enumerate(sp["rows"][:5]):
        y = 3.72 + r * 0.62
        if r % 2 == 0:
            rect(ctx, slide, 0.6, y - 0.07, 9.35, 0.56, "STRIPE")
        text(ctx, slide, 0.78, y, 3.4, 0.42, row["label"], size=12.5, valign="middle")
        for c, v in enumerate(row["cells"][:3]):
            cx = col_x[c] + col_w / 2
            if v == "check":
                image(ctx, slide, ctx.icon("check", "ACCENT"), cx - 0.13, y + 0.07, w=0.26, h=0.26)
            elif v == "times":
                image(ctx, slide, ctx.icon("times", "RED"), cx - 0.13, y + 0.07, w=0.26, h=0.26)
            elif v == "minus":
                image(ctx, slide, ctx.icon("minus", "AMBER"), cx - 0.13, y + 0.07, w=0.26, h=0.26)
            elif v == "dash":
                text(ctx, slide, col_x[c], y, col_w, 0.42, "—", size=12, color="MUT",
                     align="center", valign="middle")
            elif isinstance(v, dict):
                text(ctx, slide, col_x[c], y, col_w, 0.42, v["text"], size=12,
                     bold=True, color=v.get("color", "TXT"), align="center", valign="middle")
    if sp.get("closing_runs"):
        text(ctx, slide, 0.6, 6.95, 12.1, 0.35, sp["closing_runs"], size=13)


def s_tech_steps_loop(ctx, prs, slide, sp, meta):
    for i, p in enumerate(sp["steps"][:3]):
        y = 1.85 + i * 1.05
        icon_circle(ctx, slide, p["icon"], "SUB", 0.6, y, 0.54)
        text(ctx, slide, 1.35, y - 0.04, 5.6, 0.32, p["head"], size=14.5, bold=True, color="HEAD")
        text(ctx, slide, 1.35, y + 0.29, 5.6, 0.6, p["body"], size=12, color="MUT")
    if sp.get("image"):
        img = Image.open(ctx.asset(sp["image"]))
        iw = 5.4
        ih = iw * img.height / img.width
        image(ctx, slide, sp["image"], 7.35, 1.8, w=iw, h=ih)
        if sp.get("image_caption"):
            text(ctx, slide, 7.35, 1.8 + ih + 0.03, 5.4, 0.3, sp["image_caption"],
                 size=10.5, italic=True, color="MUT", align="center")
    loop = sp.get("loop")
    if loop:
        rect(ctx, slide, 0.6, 5.42, 12.13, 1.5, "LIGHT", radius=0.1)
        text(ctx, slide, 0.88, 5.56, 4.5, 0.3, loop["head"], size=13, bold=True, color="HEAD")
        for i, l in enumerate(loop["items"][:4]):
            x = 0.9 + i * 3.0
            icon_circle(ctx, slide, l["icon"], "SUB", x, 5.95, 0.44, bg="WHITE")
            text(ctx, slide, x + 0.54, 5.94, 2.3, 0.26, l["head"], size=12, bold=True, color="SUB")
            text(ctx, slide, x + 0.54, 6.2, 2.35, 0.55, l["body"], size=10, color="MUT")
            if i < 3:
                text(ctx, slide, x + 2.62, 5.98, 0.4, 0.3, "→", size=14, bold=True,
                     color="ACCENT")


def s_traction_evidence(ctx, prs, slide, sp, meta):
    for i, r in enumerate(sp["result_tiles"][:3]):
        y = 1.85 + i * 1.28
        rect(ctx, slide, 0.6, y, 3.55, 1.12, "LIGHT", radius=0.1)
        text(ctx, slide, 0.82, y + 0.1, 3.1, 0.52, r["big"], size=26, bold=True, color="SUB")
        text(ctx, slide, 0.82, y + 0.64, 3.15, 0.4, r["small"], size=11.5, color="MUT")
    if sp.get("tiles_footnote"):
        text(ctx, slide, 0.6, 5.75, 3.8, 0.6, sp["tiles_footnote"], size=11, color="MUT")
    pc = sp.get("pilot_card")
    if pc:
        rect(ctx, slide, 4.5, 1.85, 5.0, 4.4, "WHITE", line="BORDER", radius=0.12)
        img = Image.open(ctx.asset(pc["image"]))
        iw = 4.5
        ih = iw * img.height / img.width
        image(ctx, slide, pc["image"], 4.75, 2.15, w=iw, h=ih)
        text(ctx, slide, 4.75, 2.15 + ih + 0.18, 4.5, 0.3, pc["head"], size=13, bold=True,
             color="HEAD")
        text(ctx, slide, 4.75, 2.15 + ih + 0.51, 4.5, 1.1, pc["body"], size=11.5, color="MUT")
    ev = sp.get("evidence")
    if ev:
        text(ctx, slide, 9.9, 1.85, 2.9, 0.3, ev.get("label", "Evidence"), size=12,
             bold=True, color="MUT")
        pos = [(9.9, 2.2), (11.42, 2.2), (9.9, 4.32)]
        for (x, y), p in zip(pos, ev["images"][:3]):
            img = Image.open(ctx.asset(p))
            iw = 1.38
            ih = min(2.06, iw * img.height / img.width)
            image(ctx, slide, p, x, y, w=iw, h=ih)
        if ev.get("caption"):
            text(ctx, slide, 11.42, 4.4, 1.45, 1.2, ev["caption"], size=9.5, color="MUT")
    if sp.get("bottom_runs"):
        rect(ctx, slide, 0.6, 6.55, 12.13, 0.62, "LIGHT", radius=0.09)
        text(ctx, slide, 0.85, 6.6, 11.7, 0.52, sp["bottom_runs"], size=12.5, valign="middle")


def s_market_roadmap(ctx, prs, slide, sp, meta):
    cc = sp.get("chart_card")
    if cc:
        rect(ctx, slide, 0.6, 1.8, 4.1, 3.5, "LIGHT", radius=0.1)
        text(ctx, slide, 0.85, 1.95, 3.6, 0.5, cc["label"], size=11.5, bold=True, color="HEAD")
        bar_chart(ctx, slide, 0.85, 2.5, 3.6, 2.3, cc["categories"],
                  [("Market", cc["values"])], ["SUB"],
                  num_fmt=cc.get("number_format", "0.0"), label_size=11,
                  max_val=cc.get("max_val"), gap=80)
        if cc.get("footnote"):
            text(ctx, slide, 0.85, 4.8, 3.6, 0.55, cc["footnote"], size=9.5, color="MUT")
    for i, p in enumerate(sp["phases"][:3]):
        y = 1.8 + i * 1.22
        rect(ctx, slide, 5.0, y, 7.73, 1.08, "ACCENT_TINT" if i == 0 else "LIGHT", radius=0.1)
        icon_circle(ctx, slide, p["icon"], "SUB", 5.2, y + 0.27, 0.52, bg="WHITE")
        text(ctx, slide, 5.95, y + 0.12, 1.85, 0.3, p["tag"], size=11, bold=True, color="SUB")
        text(ctx, slide, 5.95, y + 0.4, 3.6, 0.3, p["head"], size=13.5, bold=True, color="HEAD")
        text(ctx, slide, 9.0, y + 0.14, 3.55, 0.85, p["body"], size=10.5, color="MUT",
             valign="middle")
    for i, g in enumerate(sp.get("seg_photos", [])[:3]):
        x = 0.6 + i * 1.95
        image(ctx, slide, g["image"], x, 5.6, w=1.73, h=1.3)
        text(ctx, slide, x - 0.1, 6.92, 1.93, 0.26, g["caption"], size=9.5, color="MUT",
             align="center")
    if sp.get("why_runs"):
        text(ctx, slide, 6.75, 5.7, 5.95, 1.1, sp["why_runs"], size=12.5)


def s_ask_dark(ctx, prs, slide, sp, meta):
    for i, a in enumerate(sp["cards"][:3]):
        x = 0.6 + i * 4.14
        rect(ctx, slide, x, 1.95, 3.9, 2.85, "DEEP_CARD", line="DEEP_LINE", radius=0.12)
        icon_circle(ctx, slide, a["icon"], "ACCENT", x + 0.3, 2.3, 0.62, bg="DEEP_ICON")
        text(ctx, slide, x + 0.3, 3.12, 3.3, 0.4, a["head"], size=17, bold=True, color="WHITE")
        text(ctx, slide, x + 0.3, 3.56, 3.3, 1.1, a["body"], size=12, color="ICE")
    if sp.get("outcome_runs"):
        text(ctx, slide, 0.6, 5.15, 12.1, 0.45, sp["outcome_runs"], size=14)
    if sp.get("cta_runs"):
        rect(ctx, slide, 0.6, 5.85, 12.13, 1.05, "ACCENT", radius=0.1)
        text(ctx, slide, 0.95, 5.85, 11.5, 1.05, sp["cta_runs"], size=14.5, valign="middle")


def s_team_contact(ctx, prs, slide, sp, meta):
    for i, t in enumerate(sp["members"][:4]):
        y = 1.85 + i * 1.02
        icon_circle(ctx, slide, t.get("icon", "users"), "SUB", 0.6, y, 0.5)
        text(ctx, slide, 1.3, y - 0.04, 5.9, 0.32, t["head"], size=14, bold=True, color="HEAD")
        text(ctx, slide, 1.3, y + 0.28, 5.9, 0.5, t["body"], size=11.5, color="MUT")
    if sp.get("left_footnote"):
        text(ctx, slide, 0.6, 5.9, 6.5, 0.75, sp["left_footnote"], size=11.5, color="MUT")
    if sp.get("photo"):
        img = Image.open(ctx.asset(sp["photo"]))
        iw = 5.25
        ih = iw * img.height / img.width
        image(ctx, slide, sp["photo"], 7.5, 1.85, w=iw, h=ih)
        if sp.get("photo_caption"):
            text(ctx, slide, 7.5, 1.85 + ih + 0.03, 5.25, 0.3, sp["photo_caption"],
                 size=10.5, italic=True, color="MUT", align="center")
    ct = sp.get("contact")
    if ct:
        rect(ctx, slide, 7.5, 5.9, 5.25, 1.15, "DEEP", radius=0.1)
        tx = 7.8
        if ct.get("logo"):
            image(ctx, slide, ct["logo"], 7.8, 6.22, w=1.28)
            tx = 9.3
        runs = []
        for i, ln in enumerate(ct["lines"]):
            runs.append({"text": ln, "bold": i == 0,
                         "color": "WHITE" if i == 0 else "ICE",
                         "break": i < len(ct["lines"]) - 1})
        text(ctx, slide, tx, 5.9, 12.65 - tx, 1.15, runs, size=12, valign="middle")


BUILDERS = {
    "cover_dark": (s_cover_dark, True),
    "stat_tiles": (s_stat_tiles, False),
    "product_hero": (s_product_hero, False),
    "icon_rows_photo": (s_icon_rows_photo, False),
    "comparison_table": (s_comparison_table, False),
    "tech_steps_loop": (s_tech_steps_loop, False),
    "traction_evidence": (s_traction_evidence, False),
    "market_roadmap": (s_market_roadmap, False),
    "ask_dark": (s_ask_dark, True),
    "team_contact": (s_team_contact, False),
}


def build(spec_path, out_override=None):
    spec_path = os.path.abspath(spec_path)
    with open(spec_path, encoding="utf-8") as f:
        spec = json.load(f)
    ctx = Ctx(spec, os.path.dirname(spec_path))
    meta = spec.get("meta", {})
    brand = meta.get("brand", meta.get("company", "")).upper()

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    slides = spec["slides"]
    total = len(slides)
    for n, sp in enumerate(slides, start=1):
        typ = sp["type"]
        if typ not in BUILDERS:
            raise SystemExit(f"unknown slide type '{typ}' (see references/slide-library.md)")
        builder, dark = BUILDERS[typ]
        slide = prs.slides.add_slide(blank)
        rect(ctx, slide, 0, 0, SLIDE_W, SLIDE_H, "DEEP" if dark else "WHITE")
        if typ != "cover_dark":
            k = sp.get("kicker", {})
            kicker(ctx, slide, k.get("num", f"{n - 1:02d}"), k.get("label", ""), brand, dark)
            title(ctx, slide, sp["title"], dark)
            page_no(ctx, slide, n, total, dark)
        builder(ctx, prs, slide, sp, meta)
        if sp.get("notes"):
            slide.notes_slide.notes_text_frame.text = sp["notes"]

    out = out_override or spec.get("output") or os.path.splitext(spec_path)[0] + ".pptx"
    if not os.path.isabs(out):
        out = os.path.join(os.path.dirname(spec_path), out)
    prs.save(out)
    print(f"written: {out}  ({total} slides, theme={spec.get('theme', 'ocean')})")
    return out


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("spec")
    ap.add_argument("--out")
    args = ap.parse_args()
    build(args.spec, args.out)
